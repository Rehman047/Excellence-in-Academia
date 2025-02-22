import PyPDF2
from pptx import Presentation
from docx import Document
import os

class TextProcessor:
    @staticmethod
    def extract_from_pdf(file_path):
        """Extract text from PDF files"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
        return text

    @staticmethod
    def extract_from_pptx(file_path):
        """Extract text from PowerPoint files"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"Error processing PPTX: {str(e)}")
        return text

    @staticmethod
    def extract_from_docx(file_path):
        """Extract text from Word documents"""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            raise Exception(f"Error processing DOCX: {str(e)}")
        return text

    @staticmethod
    def extract_from_txt(file_path):
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"Error processing TXT: {str(e)}")

    @classmethod
    def extract_text(cls, file_path):
        """Extract text from various file formats"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = file_path.lower().split('.')[-1]
        
        extractors = {
            'pdf': cls.extract_from_pdf,
            'pptx': cls.extract_from_pptx,
            'docx': cls.extract_from_docx,
            'txt': cls.extract_from_txt
        }
        
        extractor = extractors.get(file_extension)
        if not extractor:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        return extractor(file_path)

    @staticmethod
    def preprocess_text(text):
        """Preprocess extracted text for better question generation"""
        # Remove extra whitespace and normalize line endings
        text = ' '.join(text.split())
        
        # Basic text cleaning
        text = text.replace('\t', ' ')
        text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
        
        return text
